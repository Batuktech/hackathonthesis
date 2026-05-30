from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "thesis-orchestral-defense-mn.pptx"
INDEX_PATH = OUT_DIR / "index.html"

SLIDE_W = Inches(16)
SLIDE_H = Inches(9)

BG = RGBColor(248, 246, 241)
INK = RGBColor(29, 36, 40)
MUTED = RGBColor(88, 96, 101)
LINE = RGBColor(209, 202, 190)
TEAL = RGBColor(0, 120, 108)
BRICK = RGBColor(173, 75, 55)
GOLD = RGBColor(204, 154, 70)
DARK = RGBColor(19, 45, 49)
WHITE = RGBColor(255, 255, 255)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=28, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_title(slide, number, title, kicker="Дипломын хамгаалалтын систем"):
    add_text(slide, f"{number:02d}", 0.72, 0.46, 0.7, 0.35, 15, BRICK, True, font=FONT_HEAD)
    add_text(slide, kicker.upper(), 1.45, 0.46, 4.2, 0.35, 12, MUTED, True)
    add_text(slide, title, 0.72, 0.92, 8.8, 0.9, 31, INK, True, font=FONT_HEAD)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.85), Inches(14.55), Inches(0.018))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide):
    add_text(slide, "Их сургуулийн дипломын ажлын хамгаалалтын удирдлагын систем", 0.72, 8.34, 7.8, 0.32, 10.5, MUTED)


def add_pill(slide, text, x, y, w, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shape


def add_bullet_list(slide, items, x, y, w, h, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def add_card(slide, title, body, x, y, w, h, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.28), Inches(y + 0.28), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    add_text(slide, title, x + 0.62, y + 0.2, w - 0.9, 0.42, 16, INK, True)
    add_text(slide, body, x + 0.32, y + 0.83, w - 0.64, h - 1.0, 13.5, MUTED)


def add_arrow(slide, x1, y1, x2, y2, label):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = TEAL
    line.line.width = Pt(2.2)
    add_text(slide, label, (x1 + x2) / 2 - 0.55, y1 - 0.38, 1.25, 0.28, 10.5, TEAL, True, PP_ALIGN.CENTER)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide, DARK)
    add_text(slide, "Дипломын хамгаалалтын систем", 0.85, 0.65, 10.2, 0.72, 18, GOLD, True, font=FONT_HEAD)
    add_text(slide, "Дипломын ажлын хамгаалалтыг нэг урсгалд удирдах систем", 0.85, 1.58, 12.2, 1.55, 40, WHITE, True, font=FONT_HEAD)
    add_text(
        slide,
        "Оюутан, удирдагч багш, шүүмж багш, админ гэсэн оролцогчдын ажлыг нэг платформд холбож, файл, үнэлгээ, хуваарь, шүүмжийн явцыг ил тод болгоно.",
        0.9,
        3.55,
        10.4,
        1.15,
        20,
        RGBColor(219, 228, 226),
    )
    add_pill(slide, "Next.js", 0.9, 5.25, 1.25, TEAL)
    add_pill(slide, "Fastify", 2.35, 5.25, 1.25, BRICK)
    add_pill(slide, "Prisma", 3.8, 5.25, 1.25, GOLD)
    add_pill(slide, "PostgreSQL", 5.25, 5.25, 1.55, TEAL)
    add_text(slide, "Товч танилцуулга", 0.9, 7.65, 3.0, 0.35, 13, RGBColor(190, 204, 202), True)

    # 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 2, "Системийн зорилго нь хамгаалалтын ажлыг эмхлэх")
    add_footer(slide)
    add_card(slide, "Асуудал", "Материал илгээх, багш томилох, хамгаалалтын шат, оноо, засварын явц олон сувгаар тархахад хяналт алдагдана.", 0.85, 2.35, 4.45, 2.55, BRICK)
    add_card(slide, "Шийдэл", "Бүх оролцогч нэг системд нэвтэрч, өөрийн эрхийн хүрээнд ажиллана. Явц, файл, оноо, шүүмж нэг өгөгдлийн загварт хадгалагдана.", 5.78, 2.35, 4.45, 2.55, TEAL)
    add_card(slide, "Үр дүн", "Админ зохион байгуулалтаа, багш үнэлгээ ба шүүмжээ, оюутан материалаа нэг мөр урсгалаар гүйцэтгэх боломжтой.", 10.7, 2.35, 4.45, 2.55, GOLD)
    add_text(slide, "Гол санаа", 0.9, 5.85, 1.5, 0.35, 15, TEAL, True)
    add_text(slide, "Дипломын хамгаалалтыг “файл солилцоо” биш, төлөвтэй ажлын урсгал болгон загварчилсан.", 0.9, 6.32, 12.6, 0.7, 25, INK, True, font=FONT_HEAD)

    # 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 3, "Хэрэглэгчийн дүр бүр тодорхой үүрэгтэй")
    add_footer(slide)
    roles = [
        ("Админ", "Хэрэглэгч, улирал, зэрэг, шүүмжийн бүлэг, хамгаалалтын хуваарь, статистикийг удирдана.", BRICK),
        ("Оюутан", "Сэдэв, товч агуулга, түлхүүр үг, ажлын файлууд, хамгаалалтын оноо, шүүмжийн засвараа харна.", TEAL),
        ("Удирдагч багш", "Өөрийн удирдсан оюутнуудын ажлыг хянаж, эхний хамгаалалтын үнэлгээнд оролцоно.", GOLD),
        ("Шүүмж багш", "Бүлгийн оюутнуудад санал, шаардлагатай өөрчлөлт, эцсийн дүгнэлт өгнө.", DARK),
    ]
    for i, (title, body, color) in enumerate(roles):
        add_card(slide, title, body, 0.92 + i * 3.65, 2.45, 3.25, 3.15, color)
    add_text(slide, "Эрхийн шалгалт", 1.05, 6.45, 2.0, 0.35, 15, BRICK, True)
    add_bullet_list(slide, ["Оюутан зөвхөн өөрийн дипломын ажлыг харна.", "Багш зөвхөн удирдсан эсвэл шүүмжийн бүлэгт хамаарах ажлыг үнэлнэ.", "Админ бүх зохион байгуулалтын мэдээллийг удирдана."], 1.05, 6.9, 12.8, 1.0, 16)

    # 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 4, "Дипломын ажил шаттай ажлын урсгалаар явна")
    add_footer(slide)
    steps = [
        ("Ноорог", "Оюутан сэдэв, агуулга, файл бүрдүүлнэ"),
        ("Илгээсэн", "Удирдагч болон админ хянана"),
        ("Хамгаалалт", "1-3-р шатны оноо бүртгэнэ"),
        ("Шүүмж", "Шүүмж багш санал, засвар өгнө"),
        ("Дууссан", "Эцсийн хамгаалалтад бэлэн болно"),
    ]
    x = 0.8
    for idx, (name, body) in enumerate(steps):
        color = [BRICK, TEAL, GOLD, DARK, TEAL][idx]
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + idx * 3.0), Inches(3.0), Inches(1.25), Inches(1.25)).fill.solid()
        oval = slide.shapes[-1]
        oval.fill.fore_color.rgb = color
        oval.line.fill.background()
        add_text(slide, str(idx + 1), x + idx * 3.0, 3.22, 1.25, 0.5, 20, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, name, x + idx * 3.0 - 0.45, 4.58, 2.1, 0.35, 16, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + idx * 3.0 - 0.7, 5.05, 2.6, 0.9, 12.5, MUTED, False, PP_ALIGN.CENTER)
        if idx < 4:
            add_arrow(slide, x + idx * 3.0 + 1.33, 3.62, x + (idx + 1) * 3.0 - 0.08, 3.62, "")
    add_text(slide, "Төлөвүүд нь Prisma enum хэлбэрээр тодорхойлогдсон тул frontend, backend, өгөгдлийн сан нэг ойлголтоор ажиллана.", 1.0, 6.75, 13.2, 0.7, 20, INK, True, font=FONT_HEAD)

    # 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 5, "Архитектур нь энгийн боловч өргөтгөхөд бэлэн")
    add_footer(slide)
    layers = [
        ("Хэрэглэгчийн тал", "Next.js App Router\nTypeScript\nTailwind CSS\nFramer Motion"),
        ("API давхарга", "Fastify\nJWT нэвтрэлт\nSwagger баримт\nZod шалгалт"),
        ("Өгөгдөл", "Prisma ORM\nPostgreSQL\nMigration\nSeed өгөгдөл"),
        ("Байршуулалт", "Docker Compose\nNginx reverse proxy\nCertbot SSL\nVolume хадгалалт"),
    ]
    for idx, (title, body) in enumerate(layers):
        add_card(slide, title, body, 1.0 + idx * 3.55, 2.5, 3.1, 3.45, [TEAL, BRICK, GOLD, DARK][idx])
        if idx < 3:
            add_arrow(slide, 4.18 + idx * 3.55, 4.08, 4.48 + idx * 3.55, 4.08, "")
    add_text(slide, "Дотоод үйлчилгээний бэлтгэл", 1.0, 6.75, 3.2, 0.36, 15, TEAL, True)
    add_text(slide, "Repository-д API gateway болон model-orchestrator үйлчилгээний суурь код байгаа нь цаашид AI inference эсвэл тусдаа backend үйлчилгээ нэмэх боломжийг үлдээсэн.", 1.0, 7.16, 13.7, 0.55, 18, MUTED)

    # 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 6, "Өгөгдлийн загвар нь хамгаалалтын бодит харилцааг барина")
    add_footer(slide)
    add_text(slide, "Гол entity-үүд", 1.0, 2.32, 2.5, 0.35, 15, BRICK, True)
    add_bullet_list(slide, ["User, StudentProfile, TeacherProfile", "Thesis, ThesisFile, DefenseStage", "DefenseSchedule, DefenseScore", "CritiqueGroup, CritiqueReview", "AcademicSeason, DegreeType"], 1.0, 2.85, 5.0, 2.6, 16)
    add_text(slide, "API гадаргуу", 7.0, 2.32, 2.5, 0.35, 15, TEAL, True)
    add_bullet_list(slide, ["/api/auth/login, /api/auth/me", "/api/theses, /api/theses/:id/files", "/api/admin/users, /api/admin/schedules", "/api/critique-groups, /api/critiques", "/api/statistics/overview, /docs"], 7.0, 2.85, 6.0, 2.6, 16)
    add_text(slide, "Давуу тал", 1.0, 6.05, 1.9, 0.35, 15, GOLD, True)
    add_text(slide, "Оноо, шүүмж, файл, хуваарь тусдаа хүснэгт боловч Thesis төвтэй холбогддог. Ингэснээр хэрэглэгчийн самбар дээр нэг ажлын бүрэн түүхийг харуулах боломжтой.", 1.0, 6.48, 13.0, 0.82, 22, INK, True, font=FONT_HEAD)

    # 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 7, "Аюулгүй байдал ба байршуулалт үндсэн түвшинд шийдэгдсэн")
    add_footer(slide)
    add_card(slide, "Нэвтрэлт", "Имэйл, нууц үг, bcrypt hash, JWT token, /api/auth/me endpoint.", 1.0, 2.4, 4.1, 2.55, TEAL)
    add_card(slide, "Эрхийн хяналт", "Админ, оюутан, багшийн дүр болон багшийн төрлөөр thesis access, score permission шалгана.", 5.95, 2.4, 4.1, 2.55, BRICK)
    add_card(slide, "Байршуулалт", "Hetzner VPS дээр Docker Compose, host Nginx reverse proxy, Certbot SSL, named volumes ашиглана.", 10.9, 2.4, 4.1, 2.55, GOLD)
    add_text(slide, "Анхаарах зүйл", 1.0, 5.85, 2.4, 0.35, 15, BRICK, True)
    add_bullet_list(slide, ["Production JWT_SECRET заавал урт, санамсаргүй байх ёстой.", "Migration хийхийн өмнө postgres-data болон uploads-data volume-ийг нөөцлөх хэрэгтэй.", "Файл хадгалалт одоогоор local volume; дараагийн шатанд S3/R2 abstraction руу шилжүүлж болно."], 1.0, 6.3, 12.9, 1.25, 16)

    # 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 8, "Сүүлийн өөрчлөлтүүд deploy хийхэд чиглэсэн")
    add_footer(slide)
    add_card(slide, "CI/CD workflow", "GitHub Actions дээр backend, frontend, Docker production build шалгалт ажиллана. main руу push хийхэд Hetzner deploy job ажиллах бүтэц нэмсэн.", 0.85, 2.35, 4.45, 2.75, TEAL)
    add_card(slide, "Production compose", "PostgreSQL, backend, frontend service-үүдийг production env, healthcheck, localhost port binding, Docker volume-той болгож эмхэлсэн.", 5.78, 2.35, 4.45, 2.75, BRICK)
    add_card(slide, "Deploy script", "Server дээр exact commit checkout, compose config validation, database readiness, Prisma migration, healthcheck дарааллаар ажилладаг болгосон.", 10.7, 2.35, 4.45, 2.75, GOLD)
    add_text(slide, "Review дүгнэлт", 0.9, 6.05, 2.0, 0.35, 15, TEAL, True)
    add_text(slide, "Төсөл локал demo-оос VPS дээр давтагдахуйц deploy хийх шат руу шилжсэн. Одоо үлдэж буй гол эрсдэл нь production secret rotation, DNS тогтвортой байдал, backup discipline.", 0.9, 6.48, 13.3, 0.85, 22, INK, True, font=FONT_HEAD)

    # 9
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 9, "AI моделийн хэсэг одоогоор orchestration суурьтай")
    add_footer(slide)
    add_card(slide, "Model registry", "demo-orchestral-classifier placeholder model идэвхтэй. score analysis болон GPU inference model-ууд disabled байдлаар ирээдүйн өргөтгөлд бүртгэгдсэн.", 0.9, 2.35, 4.35, 2.95, TEAL)
    add_card(slide, "Pipeline service", "Request ирэхэд model сонгоно, disabled model-ийг хориглоно, дараа нь demo inference service рүү дамжуулна.", 5.8, 2.35, 4.35, 2.95, BRICK)
    add_card(slide, "API gateway", "Frontend request /api/inference/demo endpoint-оор орж, gateway model-orchestrator service рүү proxy хийдэг.", 10.7, 2.35, 4.35, 2.95, GOLD)
    add_text(slide, "Одоогийн хязгаар", 0.95, 6.25, 2.5, 0.35, 15, BRICK, True)
    add_bullet_list(slide, ["Бодит ML model хараахан холбогдоогүй, placeholder хариу буцаадаг.", "Input validation байгаа ч model-ийн audit log, rate limit, tenant isolation дараагийн шатанд хэрэгтэй.", "Python model server эсвэл GPU container нэмэх service boundary бэлэн байна."], 0.95, 6.7, 13.2, 1.2, 16)

    # 10
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, 10, "Hetzner deploy нь DNS, Nginx, SSL гэсэн гурван давхаргатай")
    add_footer(slide)
    layers = [
        ("DNS", "DuckDNS subdomain нь server-ийн зөв IPv4 рүү заана.\nОдоогийн server IP: 95.217.17.127"),
        ("Nginx", "80/443 traffic-ийг frontend болон /api, /health backend route руу proxy хийнэ."),
        ("Docker", "Postgres, backend, frontend container-ууд healthcheck-тэй ажиллана."),
        ("SSL", "Certbot нь public DNS зөв resolve болсон үед Let's Encrypt certificate авна."),
    ]
    for idx, (title, body) in enumerate(layers):
        add_card(slide, title, body, 0.95 + idx * 3.55, 2.45, 3.1, 3.55, [TEAL, BRICK, GOLD, DARK][idx])
    add_text(slide, "Deploy сургамж", 0.95, 6.75, 2.4, 0.35, 15, TEAL, True)
    add_text(slide, "Application layer эрүүл байсан ч DNS буруу IP дээр заавал Certbot fail болно. Иймээс эхлээд DNS -> Nginx -> container health -> HTTPS гэсэн дарааллаар шалгах нь зөв.", 0.95, 7.14, 13.2, 0.55, 18, MUTED)

    # 11
    slide = prs.slides.add_slide(blank)
    set_bg(slide, DARK)
    add_text(slide, "11", 0.82, 0.64, 0.7, 0.35, 15, GOLD, True, font=FONT_HEAD)
    add_text(slide, "Дүгнэлт ба дараагийн алхам", 0.82, 1.28, 10.5, 0.82, 34, WHITE, True, font=FONT_HEAD)
    add_text(slide, "Энэ төсөл дипломын хамгаалалтын workflow, production deploy, CI/CD pipeline, AI model orchestration суурь гэсэн дөрвөн гол чиглэлээр бүрэн системийн хэлбэрт ойртсон.", 0.86, 2.55, 12.2, 0.9, 24, RGBColor(224, 232, 230), True, font=FONT_HEAD)
    add_card(slide, "Одоо байгаа үнэ цэнэ", "Дүрд суурилсан самбар, ажлын төлөв, файл илгээх, оноо, шүүмж, статистик, Docker байршуулалт, CI/CD workflow.", 0.9, 4.3, 4.5, 2.2, TEAL)
    add_card(slide, "AI дараагийн шат", "Placeholder inference-ийг бодит model server, preprocessing, audit log, monitoring, fallback strategy-тай болгох.", 5.75, 4.3, 4.5, 2.2, BRICK)
    add_card(slide, "Production дараагийн шат", "DNS тогтворжуулах, HTTPS баталгаажуулах, secret rotation, backup automation, monitoring dashboard нэмэх.", 10.6, 4.3, 4.5, 2.2, GOLD)
    add_text(slide, "Бэлтгэсэн: төслийн repository-д үндэслэсэн Монгол хэлний товч танилцуулга", 0.9, 7.75, 8.2, 0.35, 13, RGBColor(190, 204, 202))

    prs.save(PPTX_PATH)


def build_index():
    html = f"""<!doctype html>
<html lang="mn">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Дипломын хамгаалалтын системийн танилцуулга</title>
  <style>
    :root {{
      color-scheme: light;
      --bg: #f8f6f1;
      --ink: #1d2428;
      --muted: #586065;
      --teal: #00786c;
      --brick: #ad4b37;
      --gold: #cc9a46;
      --dark: #132d31;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      min-height: 100vh;
      font-family: Aptos, "Segoe UI", Arial, sans-serif;
      background: var(--bg);
      color: var(--ink);
      display: grid;
      place-items: center;
      padding: 32px;
    }}
    main {{
      width: min(980px, 100%);
      background: #fff;
      border: 1px solid #d8d1c4;
      border-radius: 18px;
      padding: clamp(28px, 5vw, 64px);
      box-shadow: 0 24px 80px rgba(19, 45, 49, 0.12);
    }}
    .eyebrow {{
      color: var(--brick);
      font-weight: 800;
      letter-spacing: .08em;
      text-transform: uppercase;
      font-size: 13px;
      margin-bottom: 18px;
    }}
    h1 {{
      margin: 0;
      font-size: clamp(36px, 6vw, 72px);
      line-height: 1.02;
      letter-spacing: 0;
      max-width: 820px;
    }}
    p {{
      color: var(--muted);
      font-size: clamp(18px, 2vw, 22px);
      line-height: 1.55;
      max-width: 760px;
      margin: 24px 0 0;
    }}
    a {{
      display: inline-flex;
      align-items: center;
      justify-content: center;
      min-height: 52px;
      margin-top: 34px;
      padding: 0 24px;
      background: var(--teal);
      color: white;
      border-radius: 12px;
      text-decoration: none;
      font-weight: 800;
      font-size: 18px;
    }}
    .meta {{
      display: flex;
      flex-wrap: wrap;
      gap: 10px;
      margin-top: 30px;
    }}
    .meta span {{
      border: 1px solid #d8d1c4;
      border-radius: 999px;
      padding: 8px 12px;
      color: var(--dark);
      font-weight: 700;
      background: #fbfaf7;
    }}
  </style>
</head>
<body>
  <main>
    <div class="eyebrow">Монгол хэлний .pptx танилцуулга</div>
    <h1>Дипломын ажлын хамгаалалтын удирдлагын систем</h1>
    <p>
      Төслийн зорилго, хэрэглэгчийн дүрүүд, хамгаалалтын workflow, архитектур,
      өгөгдлийн загвар, AI model orchestration, CI/CD болон Hetzner deploy-ийг 11 слайдаар товч бөгөөд ойлгомжтой тайлбарласан.
    </p>
    <a href="./{PPTX_PATH.name}" download>.pptx файл татах</a>
    <div class="meta">
      <span>11 слайд</span>
      <span>Next.js + Fastify</span>
      <span>Prisma + PostgreSQL</span>
      <span>Docker Compose</span>
    </div>
  </main>
</body>
</html>
"""
    INDEX_PATH.write_text(html, encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    build_index()
    print(PPTX_PATH)
    print(INDEX_PATH)
